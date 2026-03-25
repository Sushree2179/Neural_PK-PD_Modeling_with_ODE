from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


TITLE_FONT_SIZE = Pt(30)
BULLET_FONT_SIZE = Pt(18)
TABLE_HEADER_FONT_SIZE = Pt(13)
TABLE_BODY_FONT_SIZE = Pt(11)
MAX_TABLE_DATA_ROWS = 8


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for i, bullet in enumerate(bullets):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = BULLET_FONT_SIZE


def add_image_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    caption_lines: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.7))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = TITLE_FONT_SIZE
    title_tf.paragraphs[0].font.bold = True

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.7), Inches(1.0), width=Inches(8.8))
    else:
        missing = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.8), Inches(1.0))
        missing_tf = missing.text_frame
        missing_tf.text = f"Image not found: {image_path}"
        missing_tf.paragraphs[0].font.size = BULLET_FONT_SIZE
        missing_tf.paragraphs[0].font.color.rgb = RGBColor(180, 0, 0)

    cap_box = slide.shapes.add_textbox(Inches(9.7), Inches(1.2), Inches(3.0), Inches(5.0))
    cap_tf = cap_box.text_frame
    cap_tf.clear()
    for i, line in enumerate(caption_lines):
        p = cap_tf.paragraphs[0] if i == 0 else cap_tf.add_paragraph()
        p.text = line
        p.font.size = BULLET_FONT_SIZE
        p.level = 0


def add_hyperparameter_table_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.7))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = TITLE_FONT_SIZE
    title_tf.paragraphs[0].font.bold = True

    rows = [
        ("input_dim", "2051", "Input feature width", "Capacity and memory footprint"),
        ("hidden_dim", "128", "Shared encoder width", "Representation power vs overfit risk"),
        ("latent_dim", "64", "Latent bottleneck size", "Compression vs transfer"),
        ("dropout", "0.2", "Encoder regularization", "Generalization stability"),
        ("batch_size", "64", "Batch granularity", "Gradient stability / throughput"),
        ("learning_rate", "1e-3", "Optimizer step size", "Convergence speed and stability"),
        ("weight_decay", "1e-4", "L2 regularization", "Prevents parameter blow-up"),
        ("patience", "40", "Early stopping tolerance", "Avoids premature stopping"),
        ("grad_clip", "1.0", "Gradient clipping threshold", "Prevents exploding gradients"),
        ("focal_gamma", "2.0", "Focal loss focus", "Emphasizes hard samples"),
    ]

    n_rows = len(rows) + 1
    n_cols = 4
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(1.0), Inches(12.1), Inches(5.9))
    table = table_shape.table

    headers = ["Hyperparameter", "Value", "Role", "Expected Impact"]
    for col, head in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = head
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = TABLE_HEADER_FONT_SIZE

    for i, row in enumerate(rows, start=1):
        for col, val in enumerate(row):
            cell = table.cell(i, col)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = TABLE_BODY_FONT_SIZE


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[tuple[str, ...]],
    top: float = 1.0,
    height: float = 5.9,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.7))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = TITLE_FONT_SIZE
    title_tf.paragraphs[0].font.bold = True

    # Keep tables readable by limiting rows per slide.
    if len(rows) > MAX_TABLE_DATA_ROWS:
        filler = ["..."] * len(headers)
        if len(headers) >= 3:
            filler[2] = "Continued in notes/source"
        rows = rows[: MAX_TABLE_DATA_ROWS - 1] + [tuple(filler)]

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows,
        n_cols,
        Inches(0.6),
        Inches(top),
        Inches(12.1),
        Inches(height),
    )
    table = table_shape.table

    for col, head in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = head
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = TABLE_HEADER_FONT_SIZE

    for i, row in enumerate(rows, start=1):
        for col, val in enumerate(row):
            cell = table.cell(i, col)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = TABLE_BODY_FONT_SIZE


def build_deck(output_path: Path, repo_root: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Phase 3B: Multi-Task Neural PK-PD with Neural ODE",
        "Technical deep dive | 18-slide deck",
    )

    slide_specs = [
        (
            "Agenda",
            [
                "Problem framing and motivation for a unified model.",
                "Theory behind multi-task learning and Neural ODE dynamics.",
                "Data construction and feature schema across tasks.",
                "Detailed architecture, loss design, and training protocol.",
                "Results interpretation, risks, and next experimental plan.",
            ],
        ),
        (
            "Problem Framing",
            [
                "Drug programs require simultaneous modeling of efficacy and ADMET endpoints.",
                "Single-task pipelines underuse transferable molecular information.",
                "PK-PD decisions require interpretable trajectory behavior, not only endpoint scores.",
                "Goal: combine predictive power and mechanistic consistency in one stack.",
            ],
        ),
        (
            "Why Multi-Task Learning",
            [
                "A shared encoder transfers supervision across related endpoints.",
                "Hard sharing regularizes low-sample tasks and can improve robustness.",
                "Task-specific heads preserve endpoint-specific decision boundaries.",
                "This architecture balances generalization and specialization.",
            ],
        ),
        (
            "Why Neural ODE",
            [
                "PK concentration behavior is naturally continuous-time.",
                "Neural ODE permits trajectory queries at any time grid.",
                "Dynamics can embed mechanistic priors such as first-order elimination.",
                "This improves plausibility versus purely discrete latent updates.",
            ],
        ),
        (
            "Data Pipeline Overview",
            [
                "Phase 3A processed artifacts and bridge outputs are loaded reproducibly.",
                "Task arrays are split into train, validation, and test partitions.",
                "DataLoaders are rebuilt per task with consistent batch semantics.",
                "Scaling preserves differences between regression and classification targets.",
            ],
        ),
        (
            "Feature Space Details",
            [
                "Input vector combines descriptors, fingerprints, and docking_quality.",
                "One unified input dimension is maintained for all four tasks.",
                "Structure-derived docking signal is injected into downstream learning.",
                "Consistent schema enables strict encoder sharing across tasks.",
            ],
        ),
        (
            "Structure-to-Binding-to-PK Bridge",
            [
                "RapidDock geometry quality is aggregated at protein-target level.",
                "Binding samples receive target-informed docking_quality augmentation.",
                "Other ADMET tasks are zero-padded to preserve shared schema.",
                "Bridge design injects structural signal without task-format drift.",
            ],
        ),
        (
            "Shared Encoder Architecture",
            [
                "Block 1: Linear(input_dim, hidden_dim) + LayerNorm + ReLU + Dropout.",
                "Block 2 repeats hidden transformation for richer nonlinear features.",
                "Projection layer maps hidden representation to latent_dim.",
                "LayerNorm is preferred over BatchNorm for interleaved task batches.",
            ],
        ),
        (
            "Model Hyperparameters",
            [
                "Architecture: input_dim, hidden_dim, latent_dim, dropout.",
                "Head capacity: reg_head_hidden and reg_head_dropout for binding.",
                "Optimization: learning_rate, weight_decay, batch_size, epochs, patience, grad_clip.",
                "Loss: task weights, focal_gamma, and class positive weights.",
                "Data split: test_size, val_size, and fixed random seed.",
            ],
        ),
        (
            "Why This Architecture Was Chosen",
            [
                "Shared encoder + task heads balances transfer and endpoint specialization.",
                "LayerNorm was selected for stable interleaved multi-task training.",
                "A deeper binding head addresses harder regression complexity.",
                "Logits-based heads support stable optimization and calibration.",
                "Neural ODE adds mechanistic PK trajectories, not only scalar predictions.",
            ],
        ),
        (
            "Task Heads",
            [
                "DeepRegressionHead is used for the harder binding regression task.",
                "RegressionHead is used for clearance prediction.",
                "ClassificationHead produces logits for hERG and Caco-2.",
                "Head-level asymmetry matches endpoint complexity and data profile.",
            ],
        ),
        (
            "PKODEFunc Internals",
            [
                "Latent representation predicts PK parameters CL and V.",
                "Exponentiation enforces positive-valued pharmacokinetic parameters.",
                "Elimination rate is defined as k = CL / V.",
                "Dynamics follow dC/dt = -kC and are integrated over time.",
            ],
        ),
        (
            "Training Protocol",
            [
                "Train/validation/test splits are prepared task-wise.",
                "Training interleaves task updates inside each epoch.",
                "Minimum-batch scheduling avoids repeated oversampling of small tasks.",
                "Gradient clipping is applied each optimization step for stability.",
            ],
        ),
        (
            "Loss Design",
            [
                "Objective is a weighted sum across all tasks.",
                "MSE is used for binding and clearance regression.",
                "Weighted BCE-with-logits or focal logits is used for classification.",
                "Positive-class weighting mitigates imbalance effects.",
            ],
        ),
        (
            "Optimization and Early Stopping",
            [
                "Adam optimizer is used with weight decay regularization.",
                "ReduceLROnPlateau adapts learning rate from validation behavior.",
                "Early stopping monitors aggregated validation loss with patience.",
                "Best checkpoint is persisted for reproducible handoff.",
            ],
        ),
        (
            "Evaluation Protocol",
            [
                "Regression quality is tracked with RMSE, MAE, and R2.",
                "Classification quality is tracked with AUROC, accuracy, and F1.",
                "Decision thresholds are selected on validation using Youden criterion.",
                "Task-level targets support consistent go/no-go monitoring.",
            ],
        ),
    ]

    for title, bullets in slide_specs:
        add_bullet_slide(prs, title, bullets)

    add_table_slide(
        prs,
        "Architecture Blocks Table",
        ["Module", "Inputs", "Computation", "Output", "Purpose"],
        [
            ("SharedEncoder", "features", "2x Linear+LayerNorm+ReLU+Dropout", "latent", "Shared molecular representation"),
            ("DeepRegressionHead", "latent", "2 hidden layers + dropout", "binding score", "Harder regression task"),
            ("RegressionHead", "latent", "Compact MLP", "clearance score", "Stable regression mapping"),
            ("ClassificationHead", "latent", "MLP logits", "hERG/Caco-2 logits", "Binary endpoints"),
            ("PKODEFunc", "latent + C(t)", "predict CL,V -> k; dC/dt=-kC", "PK curve", "Mechanistic PK dynamics"),
        ],
    )

    add_table_slide(
        prs,
        "Loss and Weighting Table",
        ["Task", "Loss", "Weight", "Class Handling", "Why"],
        [
            ("Binding", "MSE", "w_binding=1.8", "N/A", "Primary efficacy regression"),
            ("hERG", "BCE/Focal logits", "w_herg=1.0", "pos_weight + focal", "Safety class imbalance"),
            ("Caco-2", "BCE/Focal logits", "w_caco2=2.0", "pos_weight + focal", "Permeability emphasis"),
            ("Clearance", "MSE", "w_clearance=1.0", "N/A", "PK endpoint accuracy"),
            ("Physics", "Penalty term", "w_physics=0.1", "N/A", "Monotonic/non-negative PK curves"),
        ],
    )

    add_hyperparameter_table_slide(prs, "Hyperparameter Table Snapshot")

    add_image_slide(
        prs,
        "Training Dynamics",
        repo_root / "Coding" / "training_history.png",
        [
            "Loss curves characterize convergence and potential overfitting gap.",
            "Task metric trends identify faster and slower learning tasks.",
            "Plateau points indicate where task-specific fine-tuning is needed.",
        ],
    )

    add_image_slide(
        prs,
        "Neural ODE PK Curves",
        repo_root / "Coding" / "pk_curves.png",
        [
            "Curves are generated from continuous-time ODE integration.",
            "Monotonic decay aligns with first-order elimination assumptions.",
            "Inter-compound variation reflects inferred CL and V differences.",
        ],
    )

    add_bullet_slide(
        prs,
        "Summary and Next Work",
        [
            "Phase 3B delivers a complete and reproducible baseline pipeline.",
            "Main technical risks are task interference, imbalance sensitivity, and threshold dependence.",
            "Phase 3C will focus on targeted fine-tuning, loss-weight sweeps, and calibration.",
            "The current system is ready for controlled optimization rather than redesign.",
        ],
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main() -> None:
    repo_root = Path(__file__).resolve().parents[1]
    output_path = repo_root / "Slides" / "Phase3B_18Slide_Deck.pptx"
    build_deck(output_path, repo_root)
    print(f"Created: {output_path}")


if __name__ == "__main__":
    main()
