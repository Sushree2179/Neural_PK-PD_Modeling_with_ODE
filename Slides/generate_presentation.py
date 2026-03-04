"""
Generate Phase 2 & Phase 3 Stakeholder Deck (PPTX) — 31 slides.

Canonical generator for this repository.

Prereqs:  pip install python-pptx
Run:      python Slides/generate_presentation.py

Produces: Phase2_Phase3_Stakeholder_Deck.pptx  (16:9 widescreen)
"""

from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import copy

# ─── colour palette ───────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
MID_BLUE    = RGBColor(0x21, 0x60, 0x9B)
LIGHT_BG    = RGBColor(0xF4, 0xF6, 0xF9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MED_GRAY    = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY  = RGBColor(0xBB, 0xBB, 0xBB)
GREEN       = RGBColor(0x1E, 0x8E, 0x3E)
RED         = RGBColor(0xC0, 0x39, 0x2B)
ORANGE      = RGBColor(0xE6, 0x7E, 0x22)
TEAL        = RGBColor(0x16, 0xA0, 0x85)

# ─── layout constants (16:9 = 13.33 × 7.5 in) ────────────────────────────────
SLD_W  = Inches(13.333)
SLD_H  = Inches(7.5)
MARGIN = Inches(0.6)
TITLE_TOP    = Inches(0.3)
TITLE_H      = Inches(0.85)
SUBTITLE_TOP = Inches(1.15)
CONTENT_TOP  = Inches(1.55)
CONTENT_W    = SLD_W - 2 * MARGIN
CONTENT_H    = SLD_H - CONTENT_TOP - Inches(0.5)

OUT_FILE = Path(__file__).parent / "Phase2_Phase3_Stakeholder_Deck.pptx"


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor):
    """Set slide background to a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape_fill(slide, left, top, width, height, color: RGBColor):
    """Add a filled rectangle (used as decorative bar / banner)."""
    shape = slide.shapes.add_shape(
        1, left, top, width, height  # MSO_SHAPE.RECTANGLE = 1
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # no border
    return shape


def _tb(slide, left, top, width, height) -> "TextFrame":
    """Add a textbox and return its TextFrame."""
    return slide.shapes.add_textbox(left, top, width, height).text_frame


def _run(paragraph, text, size=20, bold=False, color=BLACK, italic=False):
    """Add a run to a paragraph (helper)."""
    r = paragraph.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


def _paragraph(tf, text, size=20, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
               space_before=Pt(6), space_after=Pt(4), bullet=False, italic=False):
    """Append a formatted paragraph to a text frame."""
    p = tf.add_paragraph() if tf.paragraphs[0].text != '' else tf.paragraphs[0]
    # If first paragraph already has text, always add_paragraph
    if p.text != '' or (len(tf.paragraphs) > 1 and tf.paragraphs[0].text != ''):
        p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    if bullet:
        p.level = 0
    _run(p, text, size=size, bold=bold, color=color, italic=italic)
    return p


# ─── slide builders ──────────────────────────────────────────────────────────

def add_title_slide(prs, title, subtitle="", date="", presenter=""):
    """Full-bleed dark title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_BLUE)
    # accent bar
    _add_shape_fill(slide, Inches(0), Inches(3.2), SLD_W, Inches(0.06), ACCENT_BLUE)
    # title
    tf = _tb(slide, MARGIN, Inches(1.0), CONTENT_W, Inches(2.0))
    tf.word_wrap = True
    _paragraph(tf, title, size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # subtitle
    if subtitle:
        tf2 = _tb(slide, MARGIN, Inches(3.55), CONTENT_W, Inches(1.5))
        tf2.word_wrap = True
        _paragraph(tf2, subtitle, size=22, color=RGBColor(0xAA, 0xCC, 0xEE))
    # footer line
    if presenter or date:
        tf3 = _tb(slide, MARGIN, Inches(6.3), CONTENT_W, Inches(0.7))
        _paragraph(tf3, f"{presenter}    |    {date}" if presenter else date,
                   size=16, color=LIGHT_GRAY)
    return slide


def add_section_slide(prs, section_title, subtitle=""):
    """Section divider (dark background, centred text)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MID_BLUE)
    _add_shape_fill(slide, Inches(0), Inches(3.5), SLD_W, Inches(0.05), ACCENT_BLUE)
    tf = _tb(slide, MARGIN, Inches(2.2), CONTENT_W, Inches(1.5))
    tf.word_wrap = True
    _paragraph(tf, section_title, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        _paragraph(tf, subtitle, size=22, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)
    return slide


def add_content_slide(prs, title, bullets=None, subbullets=None, notes=""):
    """Standard content slide with title bar + bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    # title strip
    _add_shape_fill(slide, Inches(0), Inches(0), SLD_W, Inches(1.25), DARK_BLUE)
    tf_t = _tb(slide, MARGIN, Inches(0.18), CONTENT_W, Inches(0.9))
    _paragraph(tf_t, title, size=28, bold=True, color=WHITE)

    # bullets
    if bullets:
        tf = _tb(slide, MARGIN, CONTENT_TOP, CONTENT_W, CONTENT_H)
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(10)
            p.space_after  = Pt(4)
            p.level = 0
            _run(p, f"•  {b}", size=22, color=DARK_GRAY)
            # sub-bullets
            if subbullets and i in subbullets:
                for sb in subbullets[i]:
                    sp = tf.add_paragraph()
                    sp.space_before = Pt(2)
                    sp.space_after  = Pt(2)
                    sp.level = 1
                    _run(sp, f"    ‣  {sb}", size=18, color=MED_GRAY)

    if notes:
        ns = slide.notes_slide.notes_text_frame
        ns.clear()
        ns.text = notes
    return slide


def add_table_slide(prs, title, headers, rows, notes="", col_widths=None):
    """Slide with a styled table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_shape_fill(slide, Inches(0), Inches(0), SLD_W, Inches(1.25), DARK_BLUE)
    tf_t = _tb(slide, MARGIN, Inches(0.18), CONTENT_W, Inches(0.9))
    _paragraph(tf_t, title, size=28, bold=True, color=WHITE)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_w = CONTENT_W
    tbl_h = Inches(0.45 * n_rows)
    tbl_left = MARGIN
    tbl_top = CONTENT_TOP + Inches(0.2)
    table = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h).table

    # Optional column widths
    if col_widths:
        for ci, w in enumerate(col_widths):
            table.columns[ci].width = w

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    # Data rows
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if ri % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GRAY
                p.alignment = PP_ALIGN.CENTER

    if notes:
        ns = slide.notes_slide.notes_text_frame
        ns.clear()
        ns.text = notes
    return slide


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets, notes=""):
    """Two-column layout with sub-headers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_shape_fill(slide, Inches(0), Inches(0), SLD_W, Inches(1.25), DARK_BLUE)
    tf_t = _tb(slide, MARGIN, Inches(0.18), CONTENT_W, Inches(0.9))
    _paragraph(tf_t, title, size=28, bold=True, color=WHITE)

    col_w = (CONTENT_W - Inches(0.6)) / 2

    # Left column
    _add_shape_fill(slide, MARGIN, CONTENT_TOP, col_w, Inches(0.5), ACCENT_BLUE)
    tf_lh = _tb(slide, MARGIN + Inches(0.15), CONTENT_TOP + Inches(0.05), col_w, Inches(0.4))
    _paragraph(tf_lh, left_title, size=20, bold=True, color=WHITE)

    tf_l = _tb(slide, MARGIN, CONTENT_TOP + Inches(0.6), col_w, CONTENT_H - Inches(0.6))
    tf_l.word_wrap = True
    for i, b in enumerate(left_bullets):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.space_before = Pt(8)
        _run(p, f"•  {b}", size=18, color=DARK_GRAY)

    # Right column
    right_left = MARGIN + col_w + Inches(0.6)
    _add_shape_fill(slide, right_left, CONTENT_TOP, col_w, Inches(0.5), ACCENT_BLUE)
    tf_rh = _tb(slide, right_left + Inches(0.15), CONTENT_TOP + Inches(0.05), col_w, Inches(0.4))
    _paragraph(tf_rh, right_title, size=20, bold=True, color=WHITE)

    tf_r = _tb(slide, right_left, CONTENT_TOP + Inches(0.6), col_w, CONTENT_H - Inches(0.6))
    tf_r.word_wrap = True
    for i, b in enumerate(right_bullets):
        p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
        p.space_before = Pt(8)
        _run(p, f"•  {b}", size=18, color=DARK_GRAY)

    if notes:
        ns = slide.notes_slide.notes_text_frame
        ns.clear()
        ns.text = notes
    return slide


def set_notes(slide, notes):
    ns = slide.notes_slide.notes_text_frame
    ns.clear()
    ns.text = notes


# ═══════════════════════════════════════════════════════════════════════════════
# DECK CONTENT
# ═══════════════════════════════════════════════════════════════════════════════

def build_deck():
    prs = Presentation()
    prs.slide_width  = SLD_W
    prs.slide_height = SLD_H

    # ── 1. Title ──────────────────────────────────────────────────────────────
    s = add_title_slide(
        prs,
        title="Phase 2 & Phase 3 Update:\nFeature Engineering → Multi-Task Neural ODE",
        subtitle="Neural PK–PD Modeling with Physics-Informed Neural ODE  •  PK–PD + ADMET",
        date="Snapshot: March 4, 2026",
        presenter="Subrat  •  Sushree2179/Neural_PK-PD_Modeling_with_ODE",
    )
    set_notes(s, (
        "Today I'm presenting Phase 2 and Phase 3 of the Neural PK–PD Modeling project. "
        "Phase 2 focuses on feature engineering and building a consistent multi-task dataset. "
        "Phase 3 focuses on the Neural ODE multi-task model architecture, training setup, baseline performance, "
        "and the concrete next steps to reach target metrics."
    ))

    # ── 2. Executive Summary ──────────────────────────────────────────────────
    add_content_slide(prs,
        title="Executive Summary",
        bullets=[
            "Phase 2→3 handoff complete via processed artifact (single canonical feature source)",
            "Phase 3 active: multi-variant benchmarking executed on shared encoder + Neural ODE pipeline",
            "Current feature space: 258 dims (2 physico + 256 fingerprint)",
            "Latest benchmark: clearance near target; binding / hERG / Caco-2 still below targets",
            "Next action: task-specific fine-tuning for hERG/Caco-2 with frozen shared encoder",
        ],
        notes=(
            "Phase 2→3 data handoff is stabilized around one processed artifact, and the Phase 3 pipeline is now reproducible "
            "for controlled benchmark comparisons. Latest runs keep clearance near target while binding, hERG, and Caco-2 "
            "remain below thresholds. The immediate next step is targeted fine-tuning of hERG/Caco-2 while freezing the "
            "shared encoder to improve classification performance without destabilizing the full model."
        ),
    )

    # ── 3. Scope ──────────────────────────────────────────────────────────────
    add_two_column_slide(prs,
        title="What This Review Covers",
        left_title="COVERED",
        left_bullets=[
            "Phase 2: feature engineering + dataset readiness",
            "Phase 3: model architecture + training design",
            "Baseline results + improvement plan",
        ],
        right_title="NOT COVERED",
        right_bullets=[
            "Deep Phase 1 EDA details",
            "Phase 4 deployment implementation",
        ],
        notes=(
            "This deck focuses on what reviewers and stakeholders need: the Phase 2 → Phase 3 pipeline, "
            "baseline results, and a controlled improvement plan. Phase 1 EDA exists as evidence for data quality."
        ),
    )

    # ── 4. Motivation ─────────────────────────────────────────────────────────
    add_content_slide(prs,
        title="Motivation: Why This Matters",
        bullets=[
            "One framework to predict multiple endpoints (ADMET + PK/PD-relevant outcomes)",
            "Shared representation reduces duplicated per-endpoint modeling",
            "Iterative pipeline: data → features → model → evaluation → improvement",
            "Extensible: once the pipeline is stable, improvements transfer across tasks",
        ],
        notes=(
            "The purpose is to learn a single representation that supports multiple endpoints. "
            "This helps scale modeling work and makes improvements transferable across tasks."
        ),
    )

    # ── 5. Project Phases ─────────────────────────────────────────────────────
    add_content_slide(prs,
        title="Project Context — Phases",
        bullets=[
            "Phase 1: Data exploration & validation  ✅",
            "Phase 2: Feature engineering & task-ready dataset  ✅",
            "Phase 3: Neural ODE model development  🔬  (in progress)",
            "Phase 4: Deployment  (pending)",
        ],
        notes=(
            "Phase 2 is complete and provides the model-ready dataset. Phase 3 has produced a baseline run, "
            "and our next iterations are focused on performance improvements before moving to deployment."
        ),
    )

    # ══════════════════════════════════════════════════════════════════════════
    # PHASE 2 SECTION
    # ══════════════════════════════════════════════════════════════════════════
    add_section_slide(prs, "Phase 2", "Feature Engineering & Multi-Task Dataset")

    # ── 7. Phase 2 Goals ──────────────────────────────────────────────────────
    add_content_slide(prs,
        title="Phase 2 Goals",
        bullets=[
            "Standardize inputs across sources into a single ML-ready format",
            "Ensure correct task definitions (regression vs classification)",
            "Normalize features and regression targets for stable multi-task training",
            "Prepare clean inputs for Phase 3 Neural ODE model consumption",
        ],
        notes=(
            "Phase 2 success means: training is possible without schema issues, task types are correct, "
            "and preprocessing is consistent and reproducible."
        ),
    )

    # ── 8. Training Set ───────────────────────────────────────────────────────
    add_content_slide(prs,
        title="Phase 2 Output: Final Training Set",
        bullets=[
            "Total training samples: 13,030",
            "4 prediction tasks:",
        ],
        subbullets={
            1: [
                "Binding affinity  (regression)",
                "hERG inhibition  (classification)",
                "Caco-2 permeability  (regression)",
                "Hepatocyte clearance  (regression)",
            ],
        },
        notes=(
            "The combined raw inspection totals ~347,896 records from PubChem, ChEMBL, TDC, ToxCast, PK-DB. "
            "From these, we construct 13,030 curated task training samples."
        ),
    )

    # ── 9. Feature Design ─────────────────────────────────────────────────────
    add_content_slide(prs,
        title="Phase 2 Features — What the Model Sees",
        bullets=[
            "Current feature vector: 258 dimensions",
            "256-bit molecular fingerprint (ECFP-style) + 2 physicochemical inputs",
            "Processed feature matrix is reused directly in Phase 3",
            "Feature normalization: mean = 0, std = 1",
            "Regression target normalization: StandardScaler (per task)",
        ],
        notes=(
            "This slide reflects the latest active configuration used in current benchmarks: 258 total input dimensions "
            "with fingerprint + physico descriptors integrated in the processed Phase 2 artifact."
        ),
    )

    # ── 10. Key Decision ──────────────────────────────────────────────────────
    add_content_slide(prs,
        title="Phase 2 Key Decision: Caco-2 as Regression",
        bullets=[
            "Caco-2 is a continuous permeability measurement (not binary)",
            "Correct task type ensures correct model head + loss function",
            "Prevents mismatch between target type and training objective",
            "Aligns training with actual label semantics",
        ],
        notes=(
            "Treating a continuous endpoint as classification forces a threshold and mismatched loss. "
            "We corrected Caco-2 to regression, aligning training with the label distribution."
        ),
    )

    # ── 11. Deliverables ──────────────────────────────────────────────────────
    add_content_slide(prs,
        title="Phase 2 Deliverables (Artifacts)",
        bullets=[
            "Notebook: Coding/phase1_2_data_exploration.ipynb",
            "Processed artifact: Coding/data/processed/phase2_multitask_features_with_fingerprints.csv",
        ],
        subbullets={
            0: [
                "End-to-end feature engineering + dataset preparation",
                "RDKit descriptor extraction + Morgan fingerprints",
                "Multi-task feature matrix creation + normalization",
            ],
            1: [
                "Canonical Phase 3 input source",
                "Enables reproducible benchmark comparisons",
            ],
        },
        notes=(
            "These artifacts document the full Phase 2 pipeline. Reviewers can trace the feature engineering choices "
            "and verify preprocessing and task definitions."
        ),
    )

    # ══════════════════════════════════════════════════════════════════════════
    # PHASE 3 SECTION
    # ══════════════════════════════════════════════════════════════════════════
    add_section_slide(prs, "Phase 3", "Neural ODE Multi-Task Model Development")

    # ── 13. Phase 3 Goals ─────────────────────────────────────────────────────
    add_content_slide(prs,
        title="Phase 3 Objectives",
        bullets=[
            "Multi-task neural architecture: shared encoder + task-specific heads",
            "Neural ODE integration using torchdiffeq for PK dynamics",
            "Physics-informed loss incorporating mechanistic constraints",
            "Stable multi-task training + baseline results for 4 tasks",
        ],
        notes=(
            "Phase 3 establishes the multi-task Neural ODE model and produces baseline metrics. "
            "The objective now is iterative improvement to hit target thresholds."
        ),
    )

    # ── 14. Target Metrics ────────────────────────────────────────────────────
    add_table_slide(prs,
        title="Phase 3 Target Metrics",
        headers=["Task", "Metric", "Target"],
        rows=[
            ["Binding Affinity", "R²", "> 0.60"],
            ["hERG Inhibition", "AUROC", "> 0.80"],
            ["Caco-2 Permeability", "AUROC", "> 0.75"],
            ["Clearance", "RMSE", "< 1.00"],
        ],
        notes="These targets define Phase 3 completion criteria. All four must be met to exit Phase 3.",
    )

    # ── 15. Architecture ──────────────────────────────────────────────────────
    add_content_slide(prs,
        title="Phase 3 Architecture (High Level)",
        bullets=[
            "SharedEncoder: LayerNorm → Linear → ReLU (maps input → 64-dim latent)",
            "PKODEFunc: Neural ODE dynamics in latent space via torchdiffeq",
            "Task heads: RegressionHead (binding, Caco-2, clearance) + ClassificationHead (hERG)",
            "Unified wrapper: MultiTaskPKPDModel",
            "Total trainable parameters: 44,710",
        ],
        notes=(
            "The core idea is shared representation learning with specialized heads. "
            "LayerNorm avoids BatchNorm running-statistics issues when mixing tasks in interleaved training."
        ),
    )

    # ── 16. Why Neural ODE ────────────────────────────────────────────────────
    add_content_slide(prs,
        title="Why Neural ODE?",
        bullets=[
            "Continuous-depth transformation in latent space",
            "Conceptual alignment with PK/PD dynamics (state evolves over time)",
            "Structured nonlinear mapping shared across heterogeneous endpoints",
            "Can generate PK-like concentration-time curves from learned parameters",
        ],
        notes=(
            "Neural ODE supports modeling latent evolution. Even beyond time-series tasks, it can act as a flexible "
            "shared transformation consistent with PK/PD intuition."
        ),
    )

    # ── 17. Training Setup ────────────────────────────────────────────────────
    add_content_slide(prs,
        title="Training Setup (Stability-Focused)",
        bullets=[
            "Interleaved sampling (min-step style) to prevent large-task dominance",
            "Early stopping (patience = 40 epochs)",
            "Gradient clipping = 1.0 for stability",
            "LR scheduling: ReduceLROnPlateau",
            "Loss tracks benchmarked: BCE-with-logits / focal variants + threshold tuning",
        ],
        notes=(
            "This setup is designed to produce stable, repeatable training runs so we can attribute performance changes "
            "to features/model choices rather than optimization instability."
        ),
    )

    # ── 18. Configuration ─────────────────────────────────────────────────────
    add_table_slide(prs,
        title="Model Configuration",
        headers=["Parameter", "Value"],
        rows=[
            ["input_dim", "258 (2 physico + 256 fingerprint)"],
            ["hidden_dim", "128"],
            ["latent_dim", "64"],
            ["dropout", "0.2"],
            ["batch_size", "64"],
            ["epochs (max)", "300"],
            ["learning_rate", "1e-3"],
            ["weight_decay", "1e-4"],
            ["patience", "40"],
            ["grad_clip", "1.0"],
        ],
        notes="Configuration reflects controlled training setup and explicit hyperparameter choices.",
    )

    # ── 19. Training Summary ──────────────────────────────────────────────────
    add_content_slide(prs,
        title="Training Run Summary (Latest Benchmark Cycle)",
        bullets=[
            "Training executed successfully end-to-end",
            "Multiple controlled variants compared (deep-head, focal, logits + threshold tuning)",
            "Caco-2 objective aligned to classification in active benchmark track",
            "Common benchmark table maintained across runs for comparability",
            "Pipeline is reproducible and provides diagnostics for iterative improvement",
        ],
        notes=(
            "The latest cycle focused on controlled objective/loss variants and threshold behavior, with comparable "
            "evaluation outputs across runs."
        ),
    )

    # ── 20. Baseline Results ──────────────────────────────────────────────────
    add_table_slide(prs,
        title="Latest Results vs Targets (March 4, 2026)",
        headers=["Task", "Metric", "Current", "Target", "Status"],
        rows=[
            ["Binding Affinity", "R²", "-0.029", "> 0.60", "✗  Not met"],
            ["hERG Inhibition", "AUROC", "0.482", "> 0.80", "✗  Not met"],
            ["Caco-2", "AUROC", "0.518", "> 0.75", "✗  Not met"],
            ["Clearance", "RMSE", "0.969", "< 1.00", "✓  Met"],
        ],
        notes=(
            "Clearance remains near target while binding, hERG, and Caco-2 are still below required thresholds. "
            "Current direction prioritizes targeted classification-task fine-tuning."
        ),
    )

    # ── 21. Interpretation ────────────────────────────────────────────────────
    add_content_slide(prs,
        title="Interpretation: What the Baseline Tells Us",
        bullets=[
            "Pipeline is functional — training is stable (engineering risk reduced)",
            "3/4 tasks underperform → objective balancing and task-specific adaptation remain limiting",
            "Latest feature integration and reproducible handoff are in place",
            "Clearance success indicates some signal is captured under current representation",
            "Next gains likely come from targeted fine-tuning before major architecture changes",
        ],
        notes=(
            "Current evidence supports targeted per-task optimization (especially hERG/Caco-2) using the stable shared "
            "encoder as a foundation."
        ),
    )

    # ── 22. Primary Next Action ───────────────────────────────────────────────
    add_content_slide(prs,
        title="Primary Next Action: Task-Specific Fine-Tuning",
        bullets=[
            "Freeze shared encoder and Neural ODE backbone",
            "Fine-tune hERG and Caco-2 heads/task losses in focused runs",
            "Tune decision thresholds and monitor AUROC/F1 jointly",
            "Compare against latest canonical benchmark table",
            "Promote only reproducible improvements to default configuration",
        ],
        notes=(
            "Fine-tuning targeted heads is the shortest path to improving classification metrics while preserving the "
            "stability of the shared representation."
        ),
    )

    # ── 23. Secondary Actions ─────────────────────────────────────────────────
    add_content_slide(prs,
        title="Secondary Next Actions (After Focused Fine-Tuning)",
        bullets=[
            "Hyperparameter tuning: hidden_dim / latent_dim / learning rate",
            "Adjust per-task loss weights and sampling ratios",
            "For hERG/Caco-2: PR-AUC + calibration and threshold analysis",
            "Data splits & leakage checks if needed",
        ],
        notes=(
            "We avoid changing many things at once. First optimize targeted task heads; then run broader hyperparameter and "
            "loss-balancing iterations."
        ),
    )

    # ── 24. Risks ─────────────────────────────────────────────────────────────
    add_two_column_slide(prs,
        title="Risks & Mitigations",
        left_title="RISKS",
        left_bullets=[
            "Metrics remain low after fingerprint upgrade",
            "hERG imbalance / label noise limits AUROC",
            "Task interference in shared trunk",
        ],
        right_title="MITIGATIONS",
        right_bullets=[
            "Add richer descriptors or pretrained embeddings",
            "Verify splits, add PR-AUC, recalibrate weights",
            "Task adapters / partial sharing approach",
        ],
        notes=(
            "The plan is staged: representation → tuning → architecture refinement if necessary. "
            "Each mitigation is a controlled step."
        ),
    )

    # ── 25. Reviewer Ask ──────────────────────────────────────────────────────
    add_content_slide(prs,
        title="What We Need from Reviewers / Team",
        bullets=[
            "Confirm metric priorities per task (R² vs RMSE; AUROC vs PR-AUC)",
            "Agree on Phase 3 exit criteria — what thresholds enable Phase 4?",
            "Feedback: keep Neural ODE vs compare to simpler baselines first?",
        ],
        notes=(
            "Reviewer input prevents optimizing the wrong objective. "
            "We want agreement on what constitutes success before moving to deployment work."
        ),
    )

    # ── 26. Timeline ──────────────────────────────────────────────────────────
    add_content_slide(prs,
        title="Near-Term Plan (Sequencing)",
        bullets=[
            "Week 1: hERG/Caco-2 focused fine-tuning with frozen shared encoder",
            "Week 2: threshold calibration + loss-weight/sampling adjustments",
            "Week 3: consolidate results + decide Phase 3 exit / Phase 4 scope",
        ],
        notes="Dates can be adjusted; this is sequencing rather than strict scheduling.",
    )

    # ── 27. Closing ───────────────────────────────────────────────────────────
    s = add_title_slide(
        prs,
        title="Summary & Next Milestone",
        subtitle=(
            "Phase 2→3 integration complete  •  Phase 3 benchmarking established\n"
            "Clearance near target; binding/hERG/Caco-2 still below thresholds\n"
            "Next: targeted fine-tuning for hERG/Caco-2 with frozen shared encoder"
        ),
        date="March 4, 2026",
        presenter="Subrat",
    )
    set_notes(s, (
        "The project is in a good iteration-ready state. The next milestone is performance improvement "
        "with the 1024-bit Morgan fingerprint upgrade and minimal code changes."
    ))

    # ══════════════════════════════════════════════════════════════════════════
    # BACKUP SLIDES
    # ══════════════════════════════════════════════════════════════════════════
    add_section_slide(prs, "Backup Slides", "Additional detail for reviewers")

    # Backup A
    add_content_slide(prs,
        title="Backup A — Phase 3 Notebook Structure",
        bullets=[
            "Data loading (4 tasks, 13,030 samples)",
            "Feature assembly: physico descriptors + Morgan fingerprint",
            "Dataset & loaders: MultiTaskDataset with interleaved batching",
            "Model components: SharedEncoder → PKODEFunc → Task Heads",
            "Training: MultiTaskLoss (MSE + weighted BCE) + early stopping + LR scheduling",
            "Visualizations: training history + PK curve plotting",
        ],
        notes="For reviewers wanting traceability: key cells include config, feature assembly, model, loss, training, plots.",
    )

    # Backup B
    add_content_slide(prs,
        title="Backup B — Exact Code Change Checklist",
        bullets=[
            "Freeze shared encoder/ODE and unfreeze selected task heads",
            "Apply task-specific loss/threshold tuning for hERG and Caco-2",
            "Rerun in order: task-head config → dataloaders → training → validation sweep",
            "Record: AUROC/F1/threshold table + comparative benchmark snapshot",
        ],
        notes="Emphasize controlled experiments and strict comparability against the latest benchmark table.",
    )

    # Glossary
    add_content_slide(prs,
        title="Appendix: Glossary",
        bullets=[
            "Neural ODE: dh/dt = f_θ(h,t); ODE solver integrates dynamics (continuous-depth)",
            "AUROC: ranking quality across thresholds (0.5 = random → 1.0 = perfect)",
            "R²: variance explained (0 = mean predictor; < 0 = worse than mean)",
            "RMSE: error magnitude in target units (lower is better; penalizes large errors)",
            "Morgan / ECFP4: circular fingerprint (radius=2); more bits → fewer collisions → better signal",
        ],
        notes=(
            "Use this slide to align terminology. Main practical takeaway: increasing fingerprint bits reduces collisions "
            "and usually improves predictive signal."
        ),
    )

    return prs


# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = build_deck()
    prs.save(str(OUT_FILE))
    n = len(prs.slides)
    print(f"✅  Saved {n}-slide deck → {OUT_FILE}")
    print("    Open in PowerPoint / Google Slides / LibreOffice Impress.")
