from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(34)


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(30)

    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)


def add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[tuple[str, ...]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.7))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(30)
    title_tf.paragraphs[0].font.bold = True

    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.6), Inches(1.0), Inches(12.1), Inches(5.8))
    table = table_shape.table

    for col, head in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = head
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(13)

    for i, row in enumerate(rows, start=1):
        for col, val in enumerate(row):
            cell = table.cell(i, col)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(11)


def add_diagram_text_slide(prs: Presentation, title: str, mermaid_text: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.7))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(30)
    title_tf.paragraphs[0].font.bold = True

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12.1), Inches(5.9))
    body_tf = body.text_frame
    body_tf.text = "Diagram block (Mermaid):"
    body_tf.paragraphs[0].font.bold = True
    body_tf.paragraphs[0].font.size = Pt(16)

    p = body_tf.add_paragraph()
    p.text = mermaid_text
    p.font.size = Pt(12)


def add_image_slide(prs: Presentation, title: str, image_path: Path, caption: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.7))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(30)
    title_tf.paragraphs[0].font.bold = True

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.9), Inches(1.0), width=Inches(11.2))
    else:
        msg = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.2), Inches(1.0))
        msg_tf = msg.text_frame
        msg_tf.text = f"Image not found: {image_path}"
        msg_tf.paragraphs[0].font.size = Pt(16)

    cap = slide.shapes.add_textbox(Inches(0.9), Inches(6.4), Inches(11.2), Inches(0.6))
    cap_tf = cap.text_frame
    cap_tf.text = caption
    cap_tf.paragraphs[0].font.size = Pt(13)


def build_deck(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    repo_root = Path(__file__).resolve().parents[1]
    assets = repo_root / "Slides" / "assets"

    add_title_slide(
        prs,
        "RapidDock Paper-Aligned Prototype",
        "Copy-paste-ready technical deck | structure-aware PK-PD bridge",
    )

    add_bullet_slide(prs, "Motivation and Research Gap", [
        "QSAR-only approaches miss explicit 3D geometry.",
        "Docking outputs are often disconnected from downstream PK-PD models.",
        "Need an interpretable structure -> binding -> PK-PD bridge.",
    ])

    add_image_slide(
        prs,
        "End-to-End Pipeline Diagram",
        assets / "rapiddock_pipeline.png",
        "Pipeline from SMILES and PDB pocket context to PK-PD feature augmentation.",
    )

    add_table_slide(prs, "Data Sources and Assumptions", ["Component", "Source", "In Prototype", "Why"], [
        ("Ligand identity", "ChEMBL", "Real", "Real medicinal chemistry input"),
        ("Ligand 3D", "RDKit MMFF", "Generated", "Fast reproducible geometry"),
        ("Protein pocket", "RCSB PDB", "Real if mapped", "Target-aware context"),
        ("Protein features", "Synthetic", "Placeholder", "Keep prototype runnable"),
    ])

    add_bullet_slide(prs, "Distance-Biased Transformer Theory", [
        "Attention score = QK^T/sqrt(d_h) - alpha * D_joint.",
        "Distance bias encourages local geometric interactions.",
        "Model predicts ligand-ligand and ligand-protein squared distances.",
    ])

    add_table_slide(prs, "Architecture Blocks", ["Block", "Input", "Operation", "Output", "Purpose"], [
        ("Token projection", "lig/prot features", "Linear + type embed", "joint tokens", "Unify modalities"),
        ("DB-MHA", "tokens + D_joint", "distance-biased attention", "context tokens", "Geometry-aware mixing"),
        ("FFN", "context tokens", "MLP + residual", "refined tokens", "Nonlinear transformation"),
        ("Distance heads", "lig/prot token slices", "squared distance compute", "D_LL^2, D_LP^2", "Stable targets"),
    ])

    add_image_slide(
        prs,
        "Architecture Diagram",
        assets / "rapiddock_architecture.png",
        "Joint token encoder with distance-biased transformer blocks and dual distance heads.",
    )

    add_bullet_slide(prs, "Symmetry-Aware Loss", [
        "Repeated atom types create equivalent index permutations.",
        "Evaluate candidate permutations and select minimum loss.",
        "Backprop only on best permutation to respect molecular symmetry.",
    ])

    add_table_slide(prs, "Training Setup", ["Item", "Setting"], [
        ("Epochs", "50"),
        ("Optimizer", "Adam"),
        ("LR schedule", "CosineAnnealingLR (1e-3 to 1e-5)"),
        ("Gradient clipping", "1.0"),
        ("Checkpoint", "best validation loss"),
    ])

    add_bullet_slide(prs, "Reconstruction Objective (L-BFGS)", [
        "Reconstruct ligand coordinates from predicted distances and fixed protein points.",
        "Optimize coordinate set to match D_LL and D_LP under robust distance objective.",
        "Evaluate with raw RMSD and Kabsch-aligned RMSD.",
    ])

    add_table_slide(prs, "Reconstruction Metrics", ["Metric", "Meaning", "Use"], [
        ("Raw RMSD", "Direct coordinate error", "Absolute geometric fit"),
        ("Kabsch RMSD", "Alignment-invariant error", "Shape fidelity"),
        ("CDF thresholds", "Fraction below 1/2/3A", "Quality distribution"),
    ])

    add_image_slide(
        prs,
        "Training Curve (Runtime)",
        assets / "rapiddock_training_curve.png",
        "Training and validation distance-squared MSE across epochs.",
    )

    add_image_slide(
        prs,
        "Ligand Reconstruction Scatter (Runtime)",
        assets / "rapiddock_recon_scatter.png",
        "3D comparison of true ligand coordinates and reconstructed coordinates.",
    )

    add_image_slide(
        prs,
        "RMSD Histogram and CDF (Runtime)",
        assets / "rapiddock_rmsd_hist_cdf.png",
        "Distribution-level reconstruction quality using Kabsch-aligned RMSD.",
    )

    add_bullet_slide(prs, "Structure-to-Binding Bridge", [
        "docking_score = mean predicted ligand-protein distance.",
        "docking_quality = -docking_score for higher-is-better interpretation.",
        "Correlate docking score with measured pChEMBL globally and per-target.",
    ])

    add_table_slide(prs, "Bridge Dataset Schema", ["Field", "Definition", "Direction"], [
        ("target_chembl_id", "Target identifier", "categorical"),
        ("pchembl_value", "Measured affinity", "higher is better"),
        ("docking_score", "Mean predicted L-P distance", "lower is better"),
        ("docking_quality", "-docking_score", "higher is better"),
        ("real_pocket", "PDB-backed context flag", "boolean"),
    ])

    add_image_slide(
        prs,
        "Bridge Correlation Plot (Runtime)",
        assets / "rapiddock_bridge_scatter.png",
        "Docking score versus pChEMBL with per-target visual separation.",
    )

    add_image_slide(
        prs,
        "Per-Target Correlation Plot (Runtime)",
        assets / "rapiddock_per_target_corr.png",
        "Target-level Pearson correlation summary for structure-to-binding signal.",
    )

    add_image_slide(
        prs,
        "PK-PD Integration Diagram",
        assets / "rapiddock_bridge.png",
        "Bridge from predicted geometric signal to binding and downstream PK-PD model features.",
    )

    add_table_slide(prs, "Limitations and Upgrade Path", ["Limitation", "Current", "Risk", "Upgrade"], [
        ("Protein features", "Random vectors", "Missing residue signal", "ESM-2 embeddings"),
        ("Pose realism", "Conformer translation", "Geometry gap", "Vina/DiffDock"),
        ("Atom encoding", "simple one-hot type", "low chemistry richness", "richer atom features"),
        ("Scale", "lightweight", "underfitting complex regimes", "larger model/training"),
    ])

    add_bullet_slide(prs, "Key Takeaways", [
        "Prototype preserves core RapidDock geometry-learning ideas.",
        "Produces actionable bridge artifact: structure_binding_pkpd_bridge.csv.",
        "Enables structure-informed features for downstream Neural ODE PK-PD model.",
    ])

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> None:
    repo_root = Path(__file__).resolve().parents[1]
    out_path = repo_root / "Slides" / "RapidDock_Prototype_Deck.pptx"
    build_deck(out_path)
    print(f"Created: {out_path}")


if __name__ == "__main__":
    main()
