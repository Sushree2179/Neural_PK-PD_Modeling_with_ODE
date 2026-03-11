"""
Auto-generate Phase 2 & Phase 3 stakeholder deck with images + speaker notes.

Legacy/alternate generator.
Preferred canonical generator: Slides/generate_presentation.py

Prereqs:
  pip install python-pptx

Run:
    python Slides/v1_make_presentation.py

Expected images (relative paths):
    Slides/phase2_dataset_counts.png
    Slides/feature_vector_68.png
    Slides/phase2_tdc_admet.png
    Slides/phase2_chembl_binding_targets.png
    Slides/phase2_toxcast_risk_categories.png
    Slides/phase3_training_history.png
    Slides/phase3_pk_curve_demo.png
Optional:
    Slides/phase3_architecture.png
    Slides/phase3_notebook_structure.png
    Slides/phase3_change_bits_snippet.png
"""

from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# ----------------------------
# Config
# ----------------------------
BASE_DIR = Path(__file__).parent
OUT_FILE = BASE_DIR / "Phase2_Phase3_Stakeholder_Deck_v1.pptx"
IMG_DIR = BASE_DIR

# Layout tuned for 16:9 default template (13.33 x 7.5 in)
MARGIN_L = Inches(0.6)
MARGIN_R = Inches(0.6)
MARGIN_T = Inches(0.4)
MARGIN_B = Inches(0.4)

TITLE_H = Inches(0.7)
CONTENT_TOP = MARGIN_T + TITLE_H
CONTENT_H = Inches(7.5) - CONTENT_TOP - MARGIN_B
CONTENT_W = Inches(13.33) - MARGIN_L - MARGIN_R


# ----------------------------
# Helpers
# ----------------------------
def add_title(slide, title: str):
    title_box = slide.shapes.add_textbox(MARGIN_L, MARGIN_T, CONTENT_W, TITLE_H)
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = RGBColor(20, 20, 20)
    p.alignment = PP_ALIGN.LEFT


def add_bullets(slide, bullets: list[str], font_size=24):
    box = slide.shapes.add_textbox(MARGIN_L, CONTENT_TOP, CONTENT_W, CONTENT_H)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(30, 30, 30)
    return box


def add_image(slide, img_path: Path, caption: str | None = None):
    """
    Insert an image under the title, scaled to fit within content area while preserving aspect ratio.
    """
    if not img_path.exists():
        # Add a visible placeholder box if image missing
        ph = slide.shapes.add_textbox(MARGIN_L, CONTENT_TOP, CONTENT_W, CONTENT_H)
        tf = ph.text_frame
        tf.text = f"[Missing image: {img_path.as_posix()}]"
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = RGBColor(200, 0, 0)
        return

    pic = slide.shapes.add_picture(str(img_path), MARGIN_L, CONTENT_TOP, width=CONTENT_W)

    # If too tall, re-add scaled by height
    if pic.height > CONTENT_H:
        slide.shapes._spTree.remove(pic._element)
        pic = slide.shapes.add_picture(str(img_path), MARGIN_L, CONTENT_TOP, height=CONTENT_H)

    if caption:
        cap_h = Inches(0.4)
        cap_top = CONTENT_TOP + pic.height + Inches(0.1)
        if cap_top + cap_h < Inches(7.5) - MARGIN_B:
            cap = slide.shapes.add_textbox(MARGIN_L, cap_top, CONTENT_W, cap_h)
            tf = cap.text_frame
            tf.text = caption
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.italic = True
            tf.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)


def set_speaker_notes(slide, notes: str):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = notes


def make_slide(prs: Presentation, title: str, bullets=None, img=None, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_title(slide, title)

    if img:
        add_image(slide, IMG_DIR / img)
    elif bullets:
        add_bullets(slide, bullets)

    if notes:
        set_speaker_notes(slide, notes)

    return slide


# ----------------------------
# Deck Content (Slides 1–25)
# ----------------------------
SLIDES = [
    # 1
    dict(
        title="Phase 2 & Phase 3 Update: Feature Engineering → Multi-Task Neural ODE (PK–PD + ADMET)",
        bullets=[
            "Repo: Sushree2179/Neural_PK-PD_Modeling_with_ODE",
            "Presenter: Subrat",
            "Snapshot date: Mar 8, 2026",
        ],
        notes=(
            "Today I’m presenting Phase 2 and Phase 3 of the Neural PK–PD Modeling project. "
            "Phase 2 focuses on feature engineering and building a consistent multi-task dataset. "
            "Phase 3 focuses on the Neural ODE multi-task model architecture, training setup, baseline performance, "
            "and the concrete next steps to reach target metrics."
        ),
    ),
    # 2
    dict(
        title="Executive Summary (Stakeholder Version)",
        bullets=[
            "Phase 2→3 handoff complete via processed artifact (canonical feature source)",
            "Phase 3 active: Neural ODE multi-task model benchmarked with controlled variants",
            "Current benchmark: clearance near target; binding/hERG/Caco-2 below targets",
            "Next action: task-specific fine-tuning for hERG/Caco-2 with frozen shared encoder",
        ],
        notes=(
            "Phase 2→3 data handoff is stabilized around one processed artifact, and the Phase 3 pipeline is reproducible "
            "for controlled benchmark comparisons. Latest runs keep clearance near target while binding, hERG, and Caco-2 "
            "remain below thresholds. The immediate next step is targeted fine-tuning of hERG/Caco-2 while freezing "
            "the shared encoder."
        ),
    ),
    # 3
    dict(
        title="What This Review Covers (and doesn’t)",
        bullets=[
            "Covers: Phase 2 feature engineering + dataset readiness",
            "Covers: Phase 3 architecture + training setup + baseline results + improvement plan",
            "Does not cover: deep Phase 1 EDA details",
            "Does not cover: Phase 4 deployment implementation",
        ],
        notes=(
            "This deck focuses on what reviewers and stakeholders need: the Phase 2 → Phase 3 pipeline, baseline results, "
            "and a controlled improvement plan. Phase 1 EDA exists, but we only reference it as evidence for data quality."
        ),
    ),
    # 4
    dict(
        title="Motivation",
        bullets=[
            "Unified framework to predict multiple endpoints (binding + ADMET + safety + PK-relevant behavior)",
            "Shared representation reduces duplicated per-endpoint modeling",
            "Iterative loop: data → features → model → evaluate → improve",
        ],
        notes=(
            "The purpose is to learn a single representation that supports multiple endpoints. "
            "This helps scale modeling work and makes improvements transferable across tasks."
        ),
    ),
    # 5
    dict(
        title="Project Context (Phases)",
        bullets=[
            "Phase 1: Data exploration & validation ✅",
            "Phase 2: Feature engineering & task-ready dataset ✅",
            "Phase 3: Neural ODE model development 🔬 (in progress)",
            "Phase 4: Deployment (pending)",
        ],
        notes=(
            "Phase 2 is complete and provides the model-ready dataset. Phase 3 has produced a baseline run, "
            "and our next iterations are focused on performance improvements before moving to deployment."
        ),
    ),
    # 6
    dict(
        title="Phase 2 Goals",
        bullets=[
            "Standardize inputs across sources into a single ML-ready format",
            "Ensure correct task definitions (regression vs classification)",
            "Normalize features and regression targets for stable multi-task training",
            "Prepare clean inputs for Phase 3 Neural ODE model",
        ],
        notes=(
            "Phase 2 success means: training is possible without schema issues, task types are correct, "
            "and preprocessing is consistent and reproducible."
        ),
    ),
    # 7 (image)
    dict(
        title="Phase 2: Dataset Loading Summary",
        img="phase2_dataset_counts.png",
        notes=(
            "This image shows record counts loaded from each source and the combined total. "
            "It’s evidence that data access is correct and that the pipeline loads the expected datasets cleanly."
        ),
    ),
    # 8 (image)
    dict(
        title="Phase 2: Feature Vector (Baseline)",
        img="feature_vector_68.png",
        notes=(
            "Current active configuration uses 258 dimensions total (2 physico + 256 fingerprint). "
            "This slide is retained as historical baseline feature design context."
        ),
    ),
    # 9 (image)
    dict(
        title="Phase 2: TDC ADMET Distributions",
        img="phase2_tdc_admet.png",
        notes=(
            "These plots visualize ADMET endpoints: hERG class balance and continuous distributions for Caco-2 and clearance. "
            "They justify using weighted loss for hERG and treating Caco-2 as regression."
        ),
    ),
    # 10 (image)
    dict(
        title="Phase 2: ChEMBL Binding Affinity + Targets",
        img="phase2_chembl_binding_targets.png",
        notes=(
            "This slide shows the binding affinity distribution and the most represented targets. "
            "It validates that binding affinity has meaningful variance and should be learnable with sufficient features."
        ),
    ),
    # 11 (image)
    dict(
        title="Phase 2: ToxCast Safety Screening (Risk + Categories)",
        img="phase2_toxcast_risk_categories.png",
        notes=(
            "This shows safety risk stratification and category distribution from ToxCast. "
            "Even if not used as a direct training target in the current 4-task baseline, it is crucial context for "
            "future safety constraints and risk-aware modeling."
        ),
    ),
    # 12
    dict(
        title="Phase 2 Key Decision: Caco-2 as Regression",
        bullets=[
            "Historical correction: Caco-2 is continuous permeability (not binary)",
            "Current benchmark track evaluates Caco-2 as classification (AUROC) for comparability",
            "Objective alignment now tracked explicitly in experiment notes",
        ],
        notes=(
            "Caco-2 objective handling was a key experimental variable. The current benchmark table uses AUROC for Caco-2 "
            "to remain comparable with recent variant runs."
        ),
    ),
    # 13
    dict(
        title="Phase 2 Deliverables (Artifacts)",
        bullets=[
            "Notebook: Coding/phase1_2_data_exploration.ipynb (data inspection + feature engineering)",
            "Processed artifact: Coding/data/processed/phase2_multitask_features_with_fingerprints.csv",
            "Docs: Documentation/PROJECT_SUMMARY.md (main), Coding/TROUBLESHOOTING_GUIDE.md",
        ],
        notes=(
            "These artifacts document the full Phase 2 pipeline. Reviewers can trace the feature engineering choices "
            "and verify preprocessing and task definitions."
        ),
    ),
    # 14
    dict(
        title="Phase 3 Objectives",
        bullets=[
            "Multi-task neural architecture: shared encoder + task-specific heads",
            "Neural ODE integration using torchdiffeq",
            "Stable multi-task training + baseline results",
            "Targets: Binding R²>0.6, hERG AUROC>0.8, Caco-2 AUROC>0.75, Clearance RMSE<1.0",
        ],
        notes=(
            "Phase 3 establishes the multi-task Neural ODE model and produces baseline metrics. "
            "The objective now is iterative improvement to hit target thresholds."
        ),
    ),
    # 15
    dict(
        title="Phase 3 Architecture (High Level)",
        bullets=[
            "SharedEncoder (LayerNorm for interleaved multi-task stability)",
            "Task heads: regression heads + classification head (hERG)",
            "Neural ODE component (PKODEFunc) for PK-curve generation",
            "Total trainable parameters: 44,710",
        ],
        notes=(
            "The core idea is shared representation learning with specialized heads. "
            "LayerNorm avoids BatchNorm running-statistics issues when mixing tasks."
        ),
    ),
    # 16
    dict(
        title="Why Neural ODE?",
        bullets=[
            "Continuous-depth transformation in latent space",
            "Conceptual alignment with PK/PD dynamics",
            "Provides structured nonlinear mapping shared across endpoints",
        ],
        notes=(
            "Neural ODE supports modeling latent evolution. Even beyond time-series tasks, it can act as a flexible "
            "shared transformation consistent with PK/PD intuition."
        ),
    ),
    # 17
    dict(
        title="Training Setup (Stability-Focused)",
        bullets=[
            "Interleaved sampling to avoid large-task dominance",
            "Early stopping",
            "Gradient clipping = 1.0",
            "LR scheduling: ReduceLROnPlateau",
            "Loss variants benchmarked: BCE-with-logits / focal + threshold tuning",
        ],
        notes=(
            "This setup is designed to produce stable, repeatable training runs so we can attribute performance changes "
            "to features/model choices rather than optimization instability."
        ),
    ),
    # 18 (image)
    dict(
        title="Phase 3: Training History (Baseline Run)",
        img="phase3_training_history.png",
        notes=(
            "This figure shows training/validation behavior and confirms convergence and early stopping. "
            "Key point: the pipeline trains reliably; current underperformance likely stems from representation limits."
        ),
    ),
    # 19 (image)
    dict(
        title="Phase 3: Neural ODE PK Curve Demo",
        img="phase3_pk_curve_demo.png",
        notes=(
            "This plot demonstrates PK curve generation from the Neural ODE block. "
            "It validates that ODE integration works correctly and produces plausible elimination-like curves."
        ),
    ),
    # 20
    dict(
        title="Latest Results vs Targets (Mar 8, 2026)",
        bullets=[
            "Binding R²: 0.0019 (target > 0.60) → not met",
            "hERG AUROC: 0.4836 (target > 0.80) → not met",
            "Caco-2 AUROC: 0.4719 (target > 0.75) → not met",
            "Clearance RMSE: 0.9766 (target < 1.00) → met",
            "Locked thresholds: hERG=0.49, Caco-2=0.50",
        ],
        notes=(
            "Clearance remains near threshold while binding, hERG, and Caco-2 remain below targets. "
            "Current strategy prioritizes split-leakage mitigation and re-benchmarking."
        ),
    ),
    # 21
    dict(
        title="Interpretation: Likely Bottleneck",
        bullets=[
            "Training stable, but 3/4 tasks underperform → task-specific adaptation still limiting",
            "Latest feature handoff/integration is stable and reproducible",
            "Next step should fine-tune target heads before major architecture changes",
        ],
        notes=(
            "Current evidence supports targeted per-task optimization (especially hERG/Caco-2) while keeping the "
            "shared model stable."
        ),
    ),
    # 22
    dict(
        title="Primary Next Action: Split-Leakage Mitigation",
        bullets=[
            "Fix split strategy for leakage-prone tasks (starting with binding)",
            "Re-run quality gate and confirm zero overlap leakage",
            "Re-benchmark with locked thresholds and compare against Mar 8 snapshot",
        ],
        notes=(
            "Leakage control is prioritized before further model optimization so metric deltas remain trustworthy."
        ),
    ),
    # 23
    dict(
        title="Secondary Next Actions (After Focused Fine-Tuning)",
        bullets=[
            "Tune hidden_dim / latent_dim / learning rate",
            "Adjust per-task loss weights and sampling ratios",
            "For hERG/Caco-2: consider PR-AUC + calibration and threshold analysis",
        ],
        notes=(
            "We avoid changing many things at once. First optimize targeted heads; then run broader tuning."
        ),
    ),
    # 24
    dict(
        title="What We Need from Reviewers / Team",
        bullets=[
            "Confirm metric priorities per endpoint (AUROC vs PR-AUC; R² vs RMSE)",
            "Agree on Phase 3 exit criteria for Phase 4 readiness",
            "Feedback: Neural ODE value vs simpler baselines (required comparisons?)",
        ],
        notes=(
            "Reviewer input here prevents optimizing the wrong objective. "
            "We want agreement on what constitutes success before moving to deployment work."
        ),
    ),
    # 25
    dict(
        title="Appendix: Glossary",
        bullets=[
            "Neural ODE: dh/dt = fθ(h,t); ODE solver integrates dynamics (continuous-depth)",
            "AUROC: ranking quality across thresholds (0.5 random → 1.0 perfect)",
            "R²: variance explained (0 mean predictor; <0 worse than mean)",
            "RMSE: error magnitude (lower is better; penalizes large errors)",
            "Morgan/ECFP4: circular fingerprint (radius=2); more bits reduces collisions",
        ],
        notes=(
            "Use this slide to align terminology. Main practical takeaway: increasing fingerprint bits reduces collisions "
            "and usually improves predictive signal."
        ),
    ),
]


def main():
    prs = Presentation()

    # Build slides
    for i, s in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        add_title(slide, f"Slide {i} — {s['title']}")

        if s.get("img"):
            add_image(slide, IMG_DIR / s["img"])
        else:
            add_bullets(slide, s.get("bullets", []))

        set_speaker_notes(slide, s.get("notes", ""))

    prs.save(OUT_FILE)
    print(f"✅ Wrote: {OUT_FILE}")
    print("ℹ️  If any images were missing, placeholders were inserted in red.")


if __name__ == "__main__":
    main()